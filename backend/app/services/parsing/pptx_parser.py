from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.parsing.common import build_slide_asset_dir, save_binary_asset
from app.services.parsing.ocr import ocr_image_bytes
from app.utils.text_cleaning import extract_title_from_text, normalize_text


def parse_pptx(file_path: str, presentation_id: int) -> list[dict]:
    prs = Presentation(file_path)
    results: list[dict] = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        text_parts: list[str] = []
        ocr_parts: list[str] = []
        image_paths: list[str] = []
        slide_title: str | None = None

        asset_dir = build_slide_asset_dir(presentation_id, slide_number)

        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                text = normalize_text(shape.text_frame.text)
                if text:
                    text_parts.append(text)
                    if slide_title is None:
                        slide_title = extract_title_from_text(text)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_path = save_binary_asset(
                    asset_dir=asset_dir,
                    stem=f"image_{len(image_paths) + 1}",
                    ext=image.ext,
                    blob=image.blob,
                )
                image_paths.append(image_path)

                ocr_text = ocr_image_bytes(image.blob)
                if ocr_text:
                    ocr_parts.append(ocr_text)

        extracted_text = normalize_text("\n".join(text_parts))
        ocr_text = normalize_text("\n".join(ocr_parts))

        results.append(
            {
                "slide_number": slide_number,
                "slide_title": slide_title,
                "extracted_text": extracted_text,
                "ocr_text": ocr_text,
                "image_paths": image_paths,
                "slide_metadata": {
                    "source_type": "pptx",
                    "image_count": len(image_paths),
                    "text_length": len(extracted_text),
                    "ocr_length": len(ocr_text),
                },
            }
        )

    return results